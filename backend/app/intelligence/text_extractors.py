from pathlib import Path
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

def extract_text_from_file(path: Path) -> str:
    ext = path.suffix.lower()

    if ext == ".txt":
        return path.read_text(errors="ignore")

    if ext == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if ext == ".docx":
        doc = DocxDocument(str(path))
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == ".pptx":
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    raise ValueError(f"Unsupported file type: {ext}")
